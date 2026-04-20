from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
GENERATED_DIR = DOCS_DIR / "generated"


def markdown_to_pdf(markdown_path: Path, pdf_path: Path, title: str) -> None:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]

    for raw_line in markdown_path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line:
            story.append(Spacer(1, 8))
            continue
        if line.startswith("# "):
            story.append(Paragraph(line[2:], styles["Heading1"]))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], styles["Heading2"]))
        elif line.startswith("- "):
            story.append(Paragraph(f"&bull; {line[2:]}", styles["BodyText"]))
        elif line.startswith("```"):
            continue
        else:
            story.append(Paragraph(line, styles["BodyText"]))
        story.append(Spacer(1, 6))

    doc = SimpleDocTemplate(str(pdf_path), pagesize=A4)
    doc.build(story)


def build_presentation(markdown_path: Path, pptx_path: Path) -> None:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides: list[tuple[str, list[str]]] = []
    current_title = "GoldScope API"
    current_points: list[str] = []
    for raw_line in markdown_path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if line.startswith("## "):
            if current_points:
                slides.append((current_title, current_points))
            current_title = line[3:]
            current_points = []
        elif line.startswith("- "):
            current_points.append(line[2:])
    if current_points:
        slides.append((current_title, current_points))

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "GoldScope API"
    title_slide.placeholders[1].text = "Gold price intelligence platform for XJCO3011 coursework"

    for slide_title, bullet_points in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, point in enumerate(bullet_points):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = point
            paragraph.font.size = Pt(22)

    prs.save(str(pptx_path))


if __name__ == "__main__":
    api_md = DOCS_DIR / "api_documentation.md"
    report_md = DOCS_DIR / "technical_report.md"
    slides_md = DOCS_DIR / "presentation_outline.md"
    markdown_to_pdf(api_md, GENERATED_DIR / "api_documentation.pdf", "GoldScope API Documentation")
    markdown_to_pdf(report_md, GENERATED_DIR / "technical_report.pdf", "GoldScope Technical Report")
    build_presentation(slides_md, GENERATED_DIR / "goldscope_presentation.pptx")
    print("Submission assets generated successfully.")
